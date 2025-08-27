# bot.py — PDFcaster big upgrade (aiogram 3.7+)
import os
import uuid
import shutil
import subprocess
import asyncio
from typing import List
from fitz import Rect
from dotenv import load_dotenv

from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

import fitz  # PyMuPDF
import pikepdf

from aiogram import Bot, Dispatcher, F
from aiogram.types import Message, FSInputFile, ReplyKeyboardMarkup, KeyboardButton
from aiogram.enums import ParseMode
from aiogram.client.default import DefaultBotProperties
from aiogram.fsm.storage.memory import MemoryStorage
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from aiogram.filters import StateFilter

FONT_PATH = os.path.join("fonts", "arialmt.ttf")

# --- TOKEN ---
load_dotenv()  # загружаем переменные из .env
TOKEN = os.getenv("BOT_TOKEN")

bot = Bot(token=TOKEN, default=DefaultBotProperties(parse_mode=ParseMode.HTML))
dp = Dispatcher(storage=MemoryStorage())

FILES_DIR = "files"
os.makedirs(FILES_DIR, exist_ok=True)


# ----------------- STATES -----------------
class ConvertStates(StatesGroup):
    pdf2word = State()
    pdf2ppt = State()
    pdf2excel = State()
    pdf2jpg = State()
    word2pdf = State()
    ppt2pdf = State()
    excel2pdf = State()


class EditStates(StatesGroup):
    merge = State()
    split_file = State()
    split_pages = State()
    split_mode = State()
    compress = State()
    watermark_file = State()
    watermark_text = State()
    pagenum_file = State()
    pagenum_options = State()


class SecurityStates(StatesGroup):
    unlock_file = State()
    unlock_password = State()
    protect_file = State()
    protect_password = State()


# ----------------- KEYBOARDS -----------------
def main_menu():
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="Конвертировать 🔄")],
            [KeyboardButton(text="Редактировать 📝")],
            [KeyboardButton(text="Безопасность 🔐")],
        ],
        resize_keyboard=True
    )


def convert_menu():
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="PDF в Word"), KeyboardButton(text="PDF в PowerPoint")],
            [KeyboardButton(text="PDF в Excel"), KeyboardButton(text="PDF в JPG")],
            [KeyboardButton(text="Word в PDF"), KeyboardButton(text="PowerPoint в PDF")],
            [KeyboardButton(text="Excel в PDF")],
            [KeyboardButton(text="⬅ Назад")],
        ],
        resize_keyboard=True
    )


def edit_menu():
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="Объединить PDF"), KeyboardButton(text="Разделить PDF")],
            [KeyboardButton(text="Сжать PDF"), KeyboardButton(text="Водяной знак в PDF")],
            [KeyboardButton(text="Номера страниц в PDF")],
            [KeyboardButton(text="⬅ Назад")],
        ],
        resize_keyboard=True
    )


def security_menu():
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="Открыть PDF"), KeyboardButton(text="Защита PDF")],
            [KeyboardButton(text="⬅ Назад")],
        ],
        resize_keyboard=True
    )


def merge_menu():
    return ReplyKeyboardMarkup(
        keyboard=[[KeyboardButton(text="✅ Готово")], [KeyboardButton(text="⬅ Назад")]],
        resize_keyboard=True
    )


def back_menu():
    return ReplyKeyboardMarkup(
        keyboard=[[KeyboardButton(text="⬅ Назад")]],
        resize_keyboard=True
    )


# ----------------- HELPERS -----------------
def sanitize_filename(name: str) -> str:
    return "".join(c if c.isalnum() or c in " .-_()" else "_" for c in os.path.basename(name))


async def cleanup(paths: List[str]):
    for p in paths:
        try:
            if p and os.path.exists(p):
                os.remove(p)
        except Exception:
            pass


def human_size(path: str) -> str:
    s = os.path.getsize(path)
    return f"{s / 1024 / 1024:.2f} MB"


def find_bin(names):
    for n in names:
        p = shutil.which(n)
        if p:
            return p
    return None


# ----------------- START -----------------
@dp.message(F.text == "/start")
async def cmd_start(message: Message, state: FSMContext):
    await state.clear()
    await message.answer("👋 Привет! Я <b>PDFcaster_bot</b>. Что делаем?", reply_markup=main_menu())


# ----------------- NAVIGATION -----------------
@dp.message(F.text == "Конвертировать 🔄")
async def menu_convert(message: Message, state: FSMContext):
    await message.answer("Выберите конвертацию:", reply_markup=convert_menu())


@dp.message(F.text == "Редактировать 📝")
async def menu_edit(message: Message, state: FSMContext):
    await message.answer("Выберите действие редактирования:", reply_markup=edit_menu())


@dp.message(F.text == "Безопасность 🔐")
async def menu_security(message: Message, state: FSMContext):
    await message.answer("Выберите безопасность:", reply_markup=security_menu())


@dp.message(F.text == "⬅ Назад")
async def go_back(message: Message, state: FSMContext):
    await state.clear()
    await message.answer("Возврат в главное меню.", reply_markup=main_menu())


# ----------------- CONVERT HANDLERS -----------------
# PDF -> JPG
@dp.message(F.text == "PDF в JPG")
async def pdf_to_jpg_start(message: Message, state: FSMContext):
    await state.set_state(ConvertStates.pdf2jpg)
    await message.answer("Отправь PDF (будут отправлены изображения всех страниц).", reply_markup=back_menu())


@dp.message(F.document, StateFilter(ConvertStates.pdf2jpg))
async def pdf_to_jpg_file(message: Message, state: FSMContext):
    fname = sanitize_filename(message.document.file_name or f"{uuid.uuid4()}.pdf")
    input_path = os.path.join(FILES_DIR, f"pdf2jpg_{uuid.uuid4()}_{fname}")
    file_info = await bot.get_file(message.document.file_id)
    await bot.download_file(file_info.file_path, destination=input_path)

    try:
        doc = fitz.open(input_path)
        sent_files = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(alpha=False)
            img_path = os.path.join(FILES_DIR, f"{uuid.uuid4()}_page_{i+1}.png")
            pix.save(img_path)
            sent_files.append(img_path)
            await message.answer_document(FSInputFile(img_path), caption=f"Страница {i+1}")
            os.remove(img_path)
        doc.close()
        await message.answer("Готово!", reply_markup=main_menu())
    except Exception as e:
        await message.answer(f"Ошибка при конвертации PDF→JPG: {e}", reply_markup=main_menu())
    finally:
        await cleanup([input_path])
        await state.clear()


# PDF -> Word (pdf2docx)
@dp.message(F.text == "PDF в Word")
async def pdf_to_word_start(message: Message, state: FSMContext):
    await state.set_state(ConvertStates.pdf2word)
    await message.answer("Отправь PDF для конвертации в Word (.docx).", reply_markup=back_menu())


@dp.message(F.document, StateFilter(ConvertStates.pdf2word))
async def pdf_to_word_file(message: Message, state: FSMContext):
    fname = sanitize_filename(message.document.file_name or f"{uuid.uuid4()}.pdf")
    input_path = os.path.join(FILES_DIR, f"pdf2word_{uuid.uuid4()}_{fname}")
    out_docx = os.path.join(FILES_DIR, f"{os.path.splitext(fname)[0]}_{uuid.uuid4()}.docx")

    file_info = await bot.get_file(message.document.file_id)
    await bot.download_file(file_info.file_path, destination=input_path)

    try:
        try:
            from pdf2docx import Converter
        except Exception:
            await message.answer("Модуль pdf2docx не установлен. Установи: pip install pdf2docx", reply_markup=main_menu())
            return

        conv = Converter(input_path)
        conv.convert(out_docx, start=0, end=None)
        conv.close()

        await message.answer_document(FSInputFile(out_docx), caption="Готово — PDF → Word", reply_markup=main_menu())
    except Exception as e:
        await message.answer(f"Ошибка при конвертации PDF→Word: {e}", reply_markup=main_menu())
    finally:
        await cleanup([input_path, out_docx])
        await state.clear()


# PDF -> PowerPoint (pages as images → slides)
@dp.message(F.text == "PDF в PowerPoint")
async def pdf_to_ppt_start(message: Message, state: FSMContext):
    await state.set_state(ConvertStates.pdf2ppt)
    await message.answer("Отправь PDF — страницы станут слайдами PowerPoint (.pptx).", reply_markup=back_menu())


@dp.message(F.document, StateFilter(ConvertStates.pdf2ppt))
async def pdf_to_ppt_file(message: Message, state: FSMContext):
    fname = sanitize_filename(message.document.file_name or f"{uuid.uuid4()}.pdf")
    input_path = os.path.join(FILES_DIR, f"pdf2ppt_{uuid.uuid4()}_{fname}")
    out_pptx = os.path.join(FILES_DIR, f"{os.path.splitext(fname)[0]}_{uuid.uuid4()}.pptx")

    file_info = await bot.get_file(message.document.file_id)
    await bot.download_file(file_info.file_path, destination=input_path)

    try:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except Exception:
            await message.answer("Модуль python-pptx не установлен. Установи: pip install python-pptx", reply_markup=main_menu())
            return

        doc = fitz.open(input_path)
        prs = Presentation()
        for page in doc:
            pix = page.get_pixmap(alpha=False)
            img_path = os.path.join(FILES_DIR, f"{uuid.uuid4()}_ppt_page.png")
            pix.save(img_path)

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            left = top = Inches(0)
            pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)
            os.remove(img_path)

        prs.save(out_pptx)
        doc.close()

        await message.answer_document(FSInputFile(out_pptx), caption="Готово — PDF → PowerPoint", reply_markup=main_menu())
    except Exception as e:
        await message.answer(f"Ошибка при конвертации PDF→PowerPoint: {e}", reply_markup=main_menu())
    finally:
        await cleanup([input_path, out_pptx])
        await state.clear()


# PDF -> Excel (attempt camelot)
@dp.message(F.text == "PDF в Excel")
async def pdf_to_excel_start(message: Message, state: FSMContext):
    await state.set_state(ConvertStates.pdf2excel)
    await message.answer("Отправь PDF (попытаемся извлечь таблицы в Excel).", reply_markup=back_menu())


@dp.message(F.document, StateFilter(ConvertStates.pdf2excel))
async def pdf_to_excel_file(message: Message, state: FSMContext):
    fname = sanitize_filename(message.document.file_name or f"{uuid.uuid4()}.pdf")
    input_path = os.path.join(FILES_DIR, f"pdf2excel_{uuid.uuid4()}_{fname}")
    out_xlsx = os.path.join(FILES_DIR, f"{os.path.splitext(fname)[0]}_{uuid.uuid4()}.xlsx")

    file_info = await bot.get_file(message.document.file_id)
    await bot.download_file(file_info.file_path, destination=input_path)

    try:
        try:
            import camelot
        except Exception:
            await message.answer("Camelot не установлен. Установи: pip install camelot-py[cv] (требуется Java).", reply_markup=main_menu())
            return

        tables = camelot.read_pdf(input_path, pages='all')
        if not tables:
            await message.answer("Не удалось найти таблицы в PDF.", reply_markup=main_menu())
            return

        # сохраняем все таблицы в один xlsx (каждая таблица на отдельном листе)
        import pandas as pd
        with pd.ExcelWriter(out_xlsx) as writer:
            for i, table in enumerate(tables):
                df = table.df
                df.to_excel(writer, sheet_name=f"table_{i+1}", index=False)

        await message.answer_document(FSInputFile(out_xlsx), caption="Готово — PDF → Excel", reply_markup=main_menu())
    except Exception as e:
        await message.answer(f"Ошибка при конвертации PDF→Excel: {e}", reply_markup=main_menu())
    finally:
        await cleanup([input_path, out_xlsx])
        await state.clear()


# Word/PowerPoint/Excel -> PDF via libreoffice (soffice) or docx2pdf fallback for Word on Windows
def soffice_convert_to_pdf(input_path: str, out_dir: str) -> bool:
    soffice = find_bin(("soffice", "libreoffice", "soffice.exe"))
    if not soffice:
        return False
    # 'soffice --headless --convert-to pdf --outdir out_dir input_path'
    cmd = [soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, input_path]
    try:
        subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        return True
    except Exception:
        return False


@dp.message(F.text == "Word в PDF")
async def word_to_pdf_start(message: Message, state: FSMContext):
    await state.set_state(ConvertStates.word2pdf)
    await message.answer("Отправь .docx/.doc для конвертации в PDF.", reply_markup=back_menu())


@dp.message(F.document, StateFilter(ConvertStates.word2pdf))
async def word_to_pdf_file(message: Message, state: FSMContext):
    fname = sanitize_filename(message.document.file_name or f"{uuid.uuid4()}.docx")
    input_path = os.path.join(FILES_DIR, f"word2pdf_{uuid.uuid4()}_{fname}")

    file_info = await bot.get_file(message.document.file_id)
    await bot.download_file(file_info.file_path, destination=input_path)

    out_dir = FILES_DIR
    out_pdf = os.path.splitext(input_path)[0] + ".pdf"

    try:
        # try soffice
        if soffice_convert_to_pdf(input_path, out_dir):
            if os.path.exists(out_pdf):
                await message.answer_document(FSInputFile(out_pdf), caption="Готово — Word → PDF", reply_markup=main_menu())
                return
        # try docx2pdf (Windows)
        try:
            from docx2pdf import convert
            convert(input_path, out_pdf)
            if os.path.exists(out_pdf):
                await message.answer_document(FSInputFile(out_pdf), caption="Готово — Word → PDF", reply_markup=main_menu())
                return
        except Exception:
            pass

        await message.answer("Не удалось конвертировать: установите LibreOffice (soffice) или используйте Windows+MS Word (docx2pdf).", reply_markup=main_menu())
    except Exception as e:
        await message.answer(f"Ошибка при Word→PDF: {e}", reply_markup=main_menu())
    finally:
        await cleanup([input_path, out_pdf])
        await state.clear()


@dp.message(F.text == "PowerPoint в PDF")
async def ppt_to_pdf_start(message: Message, state: FSMContext):
    await state.set_state(ConvertStates.ppt2pdf)
    await message.answer("Отправь PPTX/PPT для конвертации в PDF.", reply_markup=back_menu())


@dp.message(F.document, StateFilter(ConvertStates.ppt2pdf))
async def ppt_to_pdf_file(message: Message, state: FSMContext):
    fname = sanitize_filename(message.document.file_name or f"{uuid.uuid4()}.pptx")
    input_path = os.path.join(FILES_DIR, f"ppt2pdf_{uuid.uuid4()}_{fname}")
    file_info = await bot.get_file(message.document.file_id)
    await bot.download_file(file_info.file_path, destination=input_path)

    out_pdf = os.path.splitext(input_path)[0] + ".pdf"
    try:
        if soffice_convert_to_pdf(input_path, FILES_DIR) and os.path.exists(out_pdf):
            await message.answer_document(FSInputFile(out_pdf), caption="Готово — PowerPoint → PDF", reply_markup=main_menu())
        else:
            await message.answer("Не удалось конвертировать PPT→PDF. Установи LibreOffice (soffice).", reply_markup=main_menu())
    except Exception as e:
        await message.answer(f"Ошибка при PPT→PDF: {e}", reply_markup=main_menu())
    finally:
        await cleanup([input_path, out_pdf])
        await state.clear()


@dp.message(F.text == "Excel в PDF")
async def excel_to_pdf_start(message: Message, state: FSMContext):
    await state.set_state(ConvertStates.excel2pdf)
    await message.answer("Отправь XLSX/XLS для конвертации в PDF.", reply_markup=back_menu())


@dp.message(F.document, StateFilter(ConvertStates.excel2pdf))
async def excel_to_pdf_file(message: Message, state: FSMContext):
    fname = sanitize_filename(message.document.file_name or f"{uuid.uuid4()}.xlsx")
    input_path = os.path.join(FILES_DIR, f"excel2pdf_{uuid.uuid4()}_{fname}")
    file_info = await bot.get_file(message.document.file_id)
    await bot.download_file(file_info.file_path, destination=input_path)

    out_pdf = os.path.splitext(input_path)[0] + ".pdf"
    try:
        if soffice_convert_to_pdf(input_path, FILES_DIR) and os.path.exists(out_pdf):
            await message.answer_document(FSInputFile(out_pdf), caption="Готово — Excel → PDF", reply_markup=main_menu())
        else:
            await message.answer("Не удалось конвертировать Excel→PDF. Установи LibreOffice (soffice).", reply_markup=main_menu())
    except Exception as e:
        await message.answer(f"Ошибка при Excel→PDF: {e}", reply_markup=main_menu())
    finally:
        await cleanup([input_path, out_pdf])
        await state.clear()


# ----------------- EDIT HANDLERS -----------------
# Merge
@dp.message(F.text == "Объединить PDF")
async def merge_start(message: Message, state: FSMContext):
    await state.set_state(EditStates.merge)
    await state.update_data(files=[])
    await message.answer("Отправь PDF файлы по очереди. Когда готов — нажми ✅ Готово", reply_markup=merge_menu())


@dp.message(F.document, StateFilter(EditStates.merge))
async def merge_collect(message: Message, state: FSMContext):
    data = await state.get_data()
    files = data.get("files", [])
    fname = sanitize_filename(message.document.file_name or f"{uuid.uuid4()}.pdf")
    path = os.path.join(FILES_DIR, f"merge_{uuid.uuid4()}_{fname}")
    file_info = await bot.get_file(message.document.file_id)
    await bot.download_file(file_info.file_path, destination=path)
    files.append(path)
    await state.update_data(files=files)
    await message.answer(f"Добавлен: {fname}", reply_markup=merge_menu())


@dp.message(F.text == "✅ Готово", StateFilter(EditStates.merge))
async def merge_finish(message: Message, state: FSMContext):
    data = await state.get_data()
    files = data.get("files", [])
    await state.clear()
    if len(files) < 2:
        await message.answer("Нужно минимум 2 файла.", reply_markup=main_menu())
        return
    out_path = os.path.join(FILES_DIR, f"merged_{uuid.uuid4()}.pdf")
    try:
        writer = fitz.open()
        for f in files:
            d = fitz.open(f)
            writer.insert_pdf(d)
            d.close()
        writer.save(out_path)
        writer.close()
        await message.answer_document(FSInputFile(out_path), caption="Готово — Объединено", reply_markup=main_menu())
    except Exception as e:
        await message.answer(f"Ошибка при объединении: {e}", reply_markup=main_menu())
    finally:
        await cleanup(files + [out_path])


# Split
# ----------------- SPLIT HANDLERS -----------------
# Старт разделения
@dp.message(F.text == "Разделить PDF")
async def split_start(message: Message, state: FSMContext):
    await state.set_state(EditStates.split_file)
    await message.answer("Отправь PDF для разделения.", reply_markup=back_menu())

# Получаем файл и предлагаем выбрать режим
@dp.message(F.document, StateFilter(EditStates.split_file))
async def split_receive(message: Message, state: FSMContext):
    fname = sanitize_filename(message.document.file_name or f"{uuid.uuid4()}.pdf")
    path = os.path.join(FILES_DIR, f"split_{uuid.uuid4()}_{fname}")
    file_info = await bot.get_file(message.document.file_id)
    await bot.download_file(file_info.file_path, destination=path)
    await state.update_data(file_path=path)
    await state.set_state(EditStates.split_mode)
    # Меню выбора режима
    await message.answer(
        "Выберите режим разделения PDF:",
        reply_markup=ReplyKeyboardMarkup(
            keyboard=[
                [KeyboardButton(text="Отдельные файлы")],
                [KeyboardButton(text="Один файл")],
                [KeyboardButton(text="⬅ Назад")],
            ],
            resize_keyboard=True
        )
    )

# Выбор режима
@dp.message(StateFilter(EditStates.split_mode))
async def split_mode_choice(message: Message, state: FSMContext):
    mode = message.text.strip()
    if mode not in ("Отдельные файлы", "Один файл"):
        await message.answer("Выберите один из вариантов.", reply_markup=split_mode_menu())
        return
    await state.update_data(split_mode=mode)
    await state.set_state(EditStates.split_pages)
    await message.answer("Введите номера страниц, например: 1-3,5", reply_markup=back_menu())

# Получаем страницы и выполняем разделение
@dp.message(StateFilter(EditStates.split_pages))
async def split_do(message: Message, state: FSMContext):
    data = await state.get_data()
    path = data.get("file_path")
    mode = data.get("split_mode", "Отдельные файлы")

    if not path or not os.path.exists(path):
        await message.answer("Файл не найден.", reply_markup=main_menu())
        await state.clear()
        return

    text = message.text.strip()
    pages = []
    try:
        for part in text.split(','):
            if '-' in part:
                a, b = part.split('-', 1)
                a, b = int(a), int(b)
                pages.extend(range(a - 1, b))
            else:
                pages.append(int(part) - 1)
    except Exception:
        await message.answer("Неверный формат. Используйте, например: 1-3,5", reply_markup=back_menu())
        return

    doc = fitz.open(path)
    total = len(doc)

    try:
        if mode == "Отдельные файлы":
            for p in pages:
                if 0 <= p < total:
                    out = os.path.join(FILES_DIR, f"splitpage_{uuid.uuid4()}.pdf")
                    new = fitz.open()
                    new.insert_pdf(doc, from_page=p, to_page=p)
                    new.save(out)
                    new.close()
                    await message.answer_document(FSInputFile(out), caption=f"Страница {p+1}")
                    os.remove(out)
        else:  # "Один файл"
            new = fitz.open()
            for p in pages:
                if 0 <= p < total:
                    new.insert_pdf(doc, from_page=p, to_page=p)
            out = os.path.join(FILES_DIR, f"splitpages_{uuid.uuid4()}.pdf")
            new.save(out)
            new.close()
            await message.answer_document(FSInputFile(out), caption="Выбранные страницы", reply_markup=main_menu())
            os.remove(out)

        await message.answer("Готово — разделение завершено.", reply_markup=main_menu())
    except Exception as e:
        await message.answer(f"Ошибка при разделении: {e}", reply_markup=main_menu())
    finally:
        doc.close()
        await cleanup([path])
        await state.clear()


# Compress (reliable multi-backend) — reuse logic from earlier robust handler,
# but here we include percentage reporting
@dp.message(F.text == "Сжать PDF")
async def compress_start_button(message: Message, state: FSMContext):
    await state.set_state(EditStates.compress)
    await message.answer("Отправь PDF для сжатия.", reply_markup=back_menu())


@dp.message(F.document, StateFilter(EditStates.compress))
async def compress_handler(message: Message, state: FSMContext):
    fname = sanitize_filename(message.document.file_name or f"{uuid.uuid4()}.pdf")
    input_path = os.path.join(FILES_DIR, f"compress_{uuid.uuid4()}_{fname}")
    output_path = os.path.join(FILES_DIR, f"compressed_{uuid.uuid4()}.pdf")
    file_info = await bot.get_file(message.document.file_id)
    await bot.download_file(file_info.file_path, destination=input_path)

    try:
        orig_size = os.path.getsize(input_path)
        # try Ghostscript
        gs_bin = find_bin(("gswin64c", "gswin32c", "gs"))
        if gs_bin:
            cmd = [
                gs_bin, "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.4",
                "-dPDFSETTINGS=/ebook", "-dNOPAUSE", "-dBATCH", "-dQUIET",
                f"-sOutputFile={output_path}", input_path
            ]
            try:
                subprocess.run(cmd, check=True)
                new_size = os.path.getsize(output_path)
                saved = (1 - new_size / orig_size) * 100
                await message.answer_document(FSInputFile(output_path),
                                              caption=f"✅ Сжат (Ghostscript)\nБыло: {orig_size/1024/1024:.2f} MB\nСтало: {new_size/1024/1024:.2f} MB\nЭкономия: {saved:.1f}%",
                                              reply_markup=main_menu())
                return
            except Exception:
                pass

        # try pikepdf.save() without kwargs
        try:
            with pikepdf.open(input_path) as pdf:
                pdf.save(output_path)
            if os.path.exists(output_path) and os.path.getsize(output_path) < orig_size:
                new_size = os.path.getsize(output_path)
                saved = (1 - new_size / orig_size) * 100
                await message.answer_document(FSInputFile(output_path),
                                              caption=f"✅ Сжат (pikepdf)\nБыло: {orig_size/1024/1024:.2f} MB\nСтало: {new_size/1024/1024:.2f} MB\nЭкономия: {saved:.1f}%",
                                              reply_markup=main_menu())
                return
            else:
                try:
                    os.remove(output_path)
                except Exception:
                    pass
        except Exception:
            pass

        # fallback: rasterize pages via PyMuPDF (lossy)
        try:
            doc = fitz.open(input_path)
            new_pdf = fitz.open()
            JPG_QUALITY = 70
            ZOOM = 0.9
            for page in doc:
                mat = fitz.Matrix(ZOOM, ZOOM)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img_bytes = pix.tobytes("jpeg", JPG_QUALITY)
                w, h = pix.width, pix.height
                pg = new_pdf.new_page(width=w, height=h)
                pg.insert_image(pg.rect, stream=img_bytes)
            new_pdf.save(output_path)
            new_pdf.close()
            doc.close()
            new_size = os.path.getsize(output_path)
            saved = (1 - new_size / orig_size) * 100
            await message.answer_document(FSInputFile(output_path),
                                          caption=f"✅ Сжат (растрирование)\nБыло: {orig_size/1024/1024:.2f} MB\nСтало: {new_size/1024/1024:.2f} MB\nЭкономия: {saved:.1f}%",
                                          reply_markup=main_menu())
            return
        except Exception as e:
            await message.answer(f"Не удалось сжать файл: {e}", reply_markup=main_menu())
            return

    finally:
        await cleanup([input_path, output_path])
        await state.clear()


# Watermark (text) — two-step: file -> text
# 1️⃣ Начало добавления водяного знака
@dp.message(F.text == "Водяной знак в PDF")
async def watermark_start(message: Message, state: FSMContext):
    await state.set_state(EditStates.watermark_file)
    await message.answer("Отправь PDF, куда нужно добавить текстовый водяной знак.")


# 2️⃣ Получаем PDF от пользователя
@dp.message(F.document, StateFilter(EditStates.watermark_file))
async def watermark_receive(message: Message, state: FSMContext):
    fname = message.document.file_name or f"{uuid.uuid4()}.pdf"
    path = os.path.join(FILES_DIR, f"watermark_{uuid.uuid4()}_{fname}")
    finfo = await bot.get_file(message.document.file_id)
    await bot.download_file(finfo.file_path, destination=path)
    await state.update_data(wm_file=path)
    await state.set_state(EditStates.watermark_text)
    await message.answer("Отправь текст водяного знака (короткая фраза):")


# 3️⃣ Добавляем водяной знак через ReportLab
@dp.message(StateFilter(EditStates.watermark_text))
async def watermark_text_handler(message: Message, state: FSMContext):
    data = await state.get_data()
    path = data.get("wm_file")
    if not path or not os.path.exists(path):
        await message.answer("Файл не найден.")
        await state.clear()
        return

    text = message.text.strip()
    out = os.path.join(FILES_DIR, f"watermarked_{uuid.uuid4()}.pdf")
    watermark_pdf = os.path.join(FILES_DIR, f"wm_temp_{uuid.uuid4()}.pdf")

    try:
        # Регистрируем шрифт с кириллицей
        pdfmetrics.registerFont(TTFont("CustomFont", FONT_PATH))

        # Берём размер первой страницы оригинала для водяного знака
        doc = fitz.open(path)
        page_rect = doc[0].rect
        width, height = page_rect.width, page_rect.height
        doc.close()

        # Создаём PDF с водяным знаком
        c = canvas.Canvas(watermark_pdf, pagesize=(width, height))
        c.setFont("CustomFont", 36)
        c.setFillAlpha(0.2)  # прозрачность
        cols, rows = 3, 3
        cell_width = width / cols
        cell_height = height / rows

        for i in range(rows):
            for j in range(cols):
                x = j * cell_width + cell_width / 2
                y = i * cell_height + cell_height / 2
                c.drawCentredString(x, y, text)
        c.save()

        # Накладываем водяной знак на оригинальный PDF
        doc = fitz.open(path)
        watermark = fitz.open(watermark_pdf)
        for page in doc:
            page.show_pdf_page(page.rect, watermark, 0)
        doc.save(out)
        doc.close()
        watermark.close()

        await message.answer_document(FSInputFile(out), caption="Готово — водяной знак добавлен")
    except Exception as e:
        await message.answer(f"Ошибка при добавлении водяного знака: {e}")
    finally:
        for f in [path, out, watermark_pdf]:
            if os.path.exists(f):
                os.remove(f)
        await state.clear()

# Page numbers
@dp.message(F.text == "Номера страниц в PDF")
async def pagenum_start(message: Message, state: FSMContext):
    await state.set_state(EditStates.pagenum_file)
    await message.answer("Отправь PDF, в который нужно добавить номера страниц.", reply_markup=back_menu())


@dp.message(F.document, StateFilter(EditStates.pagenum_file))
async def pagenum_receive(message: Message, state: FSMContext):
    fname = sanitize_filename(message.document.file_name or f"{uuid.uuid4()}.pdf")
    path = os.path.join(FILES_DIR, f"pnum_{uuid.uuid4()}_{fname}")
    finfo = await bot.get_file(message.document.file_id)
    await bot.download_file(finfo.file_path, destination=path)
    out = os.path.join(FILES_DIR, f"pnum_out_{uuid.uuid4()}.pdf")
    try:
        doc = fitz.open(path)
        total = len(doc)
        for i, page in enumerate(doc):
            text = f"{i+1}"
            w, h = page.rect.width, page.rect.height
            rect = Rect(0, h - 30, w, h)  # горизонтальный диапазон по ширине страницы
            page.insert_textbox(rect, text, fontsize=12, align=1)  # align=1 → по центру
        doc.save(out)
        doc.close()
        await message.answer_document(FSInputFile(out), caption="Готово — номера страниц добавлены", reply_markup=main_menu())
    except Exception as e:
        await message.answer(f"Ошибка при добавлении номеров: {e}", reply_markup=main_menu())
    finally:
        await cleanup([path, out])
        await state.clear()


# ----------------- SECURITY -----------------
# Unlock
@dp.message(F.text == "Открыть PDF")
async def unlock_start(message: Message, state: FSMContext):
    await state.set_state(SecurityStates.unlock_file)
    await message.answer("Отправь защищённый PDF.", reply_markup=back_menu())


@dp.message(F.document, StateFilter(SecurityStates.unlock_file))
async def unlock_receive(message: Message, state: FSMContext):
    fname = sanitize_filename(message.document.file_name or f"{uuid.uuid4()}.pdf")
    path = os.path.join(FILES_DIR, f"unlock_{uuid.uuid4()}_{fname}")
    finfo = await bot.get_file(message.document.file_id)
    await bot.download_file(finfo.file_path, destination=path)
    await state.update_data(unlock_file=path)
    await state.set_state(SecurityStates.unlock_password)
    await message.answer("Отправь пароль для снятия защиты:", reply_markup=back_menu())


@dp.message(StateFilter(SecurityStates.unlock_password))
async def unlock_password(message: Message, state: FSMContext):
    data = await state.get_data()
    path = data.get("unlock_file")
    pwd = message.text.strip()
    out = os.path.join(FILES_DIR, f"unlocked_{uuid.uuid4()}.pdf")
    try:
        with pikepdf.open(path, password=pwd) as pdf:
            pdf.save(out)
        await message.answer_document(FSInputFile(out), caption="Готово — пароль снят", reply_markup=main_menu())
    except Exception as e:
        await message.answer(f"Не удалось открыть PDF: {e}", reply_markup=main_menu())
    finally:
        await cleanup([path, out])
        await state.clear()


# Protect
@dp.message(F.text == "Защита PDF")
async def protect_start(message: Message, state: FSMContext):
    await state.set_state(SecurityStates.protect_file)
    await message.answer("Отправь PDF для защиты паролем.", reply_markup=back_menu())


@dp.message(F.document, StateFilter(SecurityStates.protect_file))
async def protect_receive(message: Message, state: FSMContext):
    fname = sanitize_filename(message.document.file_name or f"{uuid.uuid4()}.pdf")
    path = os.path.join(FILES_DIR, f"protect_{uuid.uuid4()}_{fname}")
    finfo = await bot.get_file(message.document.file_id)
    await bot.download_file(finfo.file_path, destination=path)
    await state.update_data(protect_file=path)
    await state.set_state(SecurityStates.protect_password)
    await message.answer("Отправь пароль, которым нужно защитить PDF:", reply_markup=back_menu())


@dp.message(StateFilter(SecurityStates.protect_password))
async def protect_password(message: Message, state: FSMContext):
    data = await state.get_data()
    path = data.get("protect_file")
    pwd = message.text.strip()
    out = os.path.join(FILES_DIR, f"protected_{uuid.uuid4()}.pdf")
    try:
        enc = pikepdf.Encryption(user=pwd, owner=pwd, R=4)
        with pikepdf.open(path) as pdf:
            pdf.save(out, encryption=enc)
        await message.answer_document(FSInputFile(out), caption="Готово — PDF защищён", reply_markup=main_menu())
    except Exception as e:
        await message.answer(f"Не удалось защитить PDF: {e}", reply_markup=main_menu())
    finally:
        await cleanup([path, out])
        await state.clear()


# ----------------- RUN -----------------
async def main():
    print("🚀 PDFcaster_bot запущен…")
    await dp.start_polling(bot)


if __name__ == "__main__":
    asyncio.run(main())
